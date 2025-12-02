# OCR还有小问题，数据库还没整合好

# ============================================================
# RAG Service (Full Version with OCR, PDF/DOCX/PPTX images,
# Enhanced Stopwords, Strong Anti-Injection, Noise Filters,
# Router, Retrieval, Reranking, Timeout-Fallback, Multi-Image,
# Multilingual Embeddings, Subject Routing)
# ============================================================

import os
import re
import shutil
import time
import json
from typing import List, Optional, Tuple
from dotenv import load_dotenv

# OCR 依赖
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\25380\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"

from PIL import Image, ImageFilter, ImageOps
from pdf2image import convert_from_path
import docx
from pptx import Presentation

# LangChain
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceBgeEmbeddings
from langchain_openai import ChatOpenAI
from langchain_core.documents import Document
from langchain_core.prompts import PromptTemplate

# 统一解析器（Parser）
from Parser import Parser, ParserError


# ============================================================
# 全局配置
# ============================================================

load_dotenv()


class SystemConfig:
    PERSIST_DIRECTORY = "chroma_db_production"

    # 使用多语言向量模型（支持中文/英文等）
    # 可以通过环境变量 EMBEDDING_MODEL_NAME 覆盖，例如：
    #   EMBEDDING_MODEL_NAME=BAAI/bge-large-zh-v1.5
    EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME", "BAAI/bge-m3")

    DEVICE = os.getenv("EMBEDDING_DEVICE", "cpu")  # 可改 "cuda"
    CHUNK_SIZE = 600
    CHUNK_OVERLAP = 120

    # 统一 OCR 语言集：Parser 和 OCRProcessorV2 都使用这个
    OCR_LANG = os.getenv("OCR_LANG", "chi_sim+eng")

    TESSERACT_TIMEOUT = 8   # OCR 超时（秒）
    RAG_TIMEOUT = 12        # GPT 生成超时（秒）


# ============================================================
# 构建停用词（中英混合 + 外部文件）
# ============================================================

def build_stopwords() -> set:
    base = {
        # 英文虚词
        "the", "is", "are", "am", "a", "an", "and", "or", "to", "for", "of", "in", "on", "at", "by", "with",
        "this", "that", "these", "those", "it", "its", "from", "as", "about", "into", "through", "during",
        "before", "after", "above", "below", "again", "further", "then", "once", "here", "there",

        # 中文常见虚词
        "的", "了", "呢", "吗", "啊", "吧", "给", "和", "以及", "如果", "然后", "因为", "所以", "请问",
        "一下", "根据", "本页", "如下", "如下所示", "例如", "比如", "如上", "如图",

        # 闲聊词（部分）
        "聊天", "天气", "今天", "明天", "你好", "请", "解释", "说明",

        # 教材无关词
        "分析", "回答", "解读", "总结", "概括",
    }

    # NLTK 英文停用词
    try:
        from nltk.corpus import stopwords
        base |= set(stopwords.words("english"))
    except Exception:
        pass

    # spaCy 停用词
    try:
        import spacy
        nlp = spacy.load("en_core_web_sm")
        base |= set(nlp.Defaults.stop_words)
    except Exception:
        pass

    # 本地 stopwords 文件
    local_path = os.path.join(os.path.dirname(__file__), "assets", "stopwords.txt")
    if os.path.exists(local_path):
        try:
            with open(local_path, "r", encoding="utf-8") as f:
                extra = {line.strip() for line in f if line.strip()}
                base |= extra
        except Exception:
            pass

    return base


GLOBAL_STOPWORDS = build_stopwords()

# 追加一批中文停用词
MORE_CN_STOPWORDS = {
    "我们", "你们", "他们", "它们",
    "这个", "那个", "这些", "那些",
    "进行", "然而", "但是", "由于", "以及", "通过", "为了", "因此",
    "例如", "比如", "主要", "一般", "通常", "必须", "需要", "如果", "同时", "所以",
    "可以看到", "如下图所示", "如下所示", "如下内容", "如下表所示",
    "一下内容", "相关内容", "更多信息",
}
GLOBAL_STOPWORDS |= MORE_CN_STOPWORDS


# ============================================================
# 学科标签 & 关键词（简单规则路由）
# ============================================================

# 学科标签固定为这 9 个
SUBJECT_LABELS = ["数", "物", "化", "生", "政", "史", "地", "计算机", "其他"]

# 简单关键词词典（中英混合，足够粗分）
SUBJECT_KEYWORDS = {
    "数": [
        "数学", "代数", "几何", "三角", "函数", "方程", "不等式",
        "微积分", "极限", "导数", "积分", "矩阵", "行列式", "线性代数",
        "概率", "统计", "数列", "随机", "数论",
        "math", "algebra", "geometry", "calculus", "probability", "statistics",
    ],
    "物": [
        "物理", "力学", "牛顿", "速度", "加速度", "匀速", "匀变速",
        "重力", "动量", "能量", "功", "功率",
        "电学", "电路", "电流", "电压", "电阻", "欧姆",
        "磁场", "光学", "折射", "反射",
        "热学", "热力学", "波动", "振动", "简谐",
        "physics", "force", "velocity", "acceleration", "circuit", "current", "voltage",
    ],
    "化": [
        "化学", "分子", "原子", "离子", "价", "化合价",
        "化学方程式", "反应", "反应热", "氧化还原", "酸碱", "盐",
        "有机", "无机", "同分异构", "烷烃", "烯烃", "芳香烃",
        "chemistry", "molecule", "atom", "reaction", "acid", "base",
    ],
    "生": [
        "生物", "细胞", "细胞膜", "细胞核", "线粒体",
        "DNA", "RNA", "基因", "染色体", "遗传", "杂交",
        "蛋白质", "酶", "代谢", "光合作用", "呼吸作用",
        "生态", "生态系统", "进化", "自然选择",
        "biology", "gene", "cell", "chromosome",
    ],
    "政": [
        "政治", "思想品德", "公民", "权利", "义务",
        "宪法", "国家", "政府", "人大", "国务院",
        "市场经济", "宏观调控", "社会主义", "中国共产党",
        "意识形态", "民主", "法治",
        "politics", "government", "constitution",
    ],
    "史": [
        "历史", "朝代", "先秦", "秦", "汉", "唐", "宋", "元", "明", "清",
        "近代史", "现代史", "世界史", "冷战", "文艺复兴",
        "革命", "战争", "工业革命",
        "history", "dynasty", "revolution",
    ],
    "地": [
        "地理", "经度", "纬度", "经纬", "地形", "盆地", "平原", "高原", "山地",
        "气候", "季风", "洋流", "降水", "气压", "锋面",
        "人口", "城市", "城市化", "区域", "版图",
        "资源", "环境", "可持续发展",
        "geography", "climate", "monsoon", "latitude", "longitude",
    ],
    "计算机": [
        "计算机", "编程", "程序", "算法", "复杂度", "数据结构",
        "链表", "栈", "队列", "树", "图", "堆",
        "操作系统", "进程", "线程", "死锁",
        "计算机网络", "TCP", "UDP", "HTTP",
        "数据库", "SQL",
        "Python", "Java", "C++", "C语言", "代码",
        "computer", "programming", "algorithm", "data structure", "network",
    ],
    "其他": [
        # 一般不放关键词，这里留空即可
    ],
}

# 支持从一些常见名字映射到学科标签
SUBJECT_SYNONYMS = {
    "数学": "数",
    "math": "数",
    "mathematics": "数",

    "物理": "物",
    "physics": "物",

    "化学": "化",
    "chemistry": "化",

    "生物": "生",
    "biology": "生",

    "政治": "政",
    "思想政治": "政",
    "思想品德": "政",
    "politics": "政",

    "历史": "史",
    "history": "史",

    "地理": "地",
    "geography": "地",

    "计算机": "计算机",
    "信息技术": "计算机",
    "computer": "计算机",
    "computer science": "计算机",
    "cs": "计算机",
}


def normalize_subject_label(raw: Optional[str]) -> Optional[str]:
    """把各种写法规整到固定 9 个学科标签之一"""
    if not raw:
        return None
    s = str(raw).strip()
    if not s:
        return None

    # 直接命中标签
    if s in SUBJECT_LABELS:
        return s

    # 同义词映射
    low = s.lower()
    if low in SUBJECT_SYNONYMS:
        return SUBJECT_SYNONYMS[low]

    if s in SUBJECT_SYNONYMS:
        return SUBJECT_SYNONYMS[s]

    # 简单 contain 匹配
    for k, v in SUBJECT_SYNONYMS.items():
        if k.lower() in low:
            return v

    return "其他"


# ============================================================
# 初始化向量库、Embedding、LLM
# ============================================================

print("[RAG] 初始化 Embedding 模型...")
embeddings = HuggingFaceBgeEmbeddings(
    model_name=SystemConfig.EMBEDDING_MODEL_NAME,
    model_kwargs={"device": SystemConfig.DEVICE},
    encode_kwargs={"normalize_embeddings": True},
)

# 关闭匿名遥测
os.environ["ANONYMIZED_TELEMETRY"] = "False"
os.environ["CHROMA_ANONYMIZED_TELEMETRY"] = "False"


def _init_vectorstore():
    persist_dir = SystemConfig.PERSIST_DIRECTORY
    try:
        return Chroma(
            persist_directory=persist_dir,
            embedding_function=embeddings,
        )
    except Exception:
        shutil.rmtree(persist_dir, ignore_errors=True)
        return Chroma(
            persist_directory=persist_dir,
            embedding_function=embeddings,
        )


vectorstore = _init_vectorstore()

# 代理 + 路由模型
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
BASE_URL = os.getenv("BASE_URL")
ROUTER_MODEL = os.getenv("ROUTER_MODEL", "gpt-4o-mini")
RAG_ANSWER_MODEL = os.getenv("RAG_ANSWER_MODEL", "gpt-4o-mini")

llm_router = ChatOpenAI(
    model=ROUTER_MODEL,
    temperature=0,
    api_key=OPENAI_API_KEY,
    base_url=BASE_URL,
)


# ============================================================
# OCR 处理模块（包含 PDF / DOCX / PPTX / PNG / JPG）
# ============================================================

class OCRProcessor:

    # -----------------------
    # 基本预处理（提高识别率）
    # -----------------------
    @staticmethod
    def preprocess_image(img: Image.Image) -> Image.Image:
        """对图片进行灰度化 / 放大 / 去噪 / 自动对比度 / 二值化，提高 OCR 识别率"""
        try:
            img = img.convert("L")  # 灰度化

            # 自动放大避免小字识别不清
            w, h = img.size
            if w < 1000:
                img = img.resize((w * 2, h * 2), Image.LANCZOS)

            # 中值滤波去噪
            img = img.filter(ImageFilter.MedianFilter(size=3))

            # 自动对比度
            img = ImageOps.autocontrast(img)

            # 二值化
            threshold = 140
            img = img.point(lambda p: 255 if p > threshold else 0)

            return img
        except Exception:
            return img

    # -----------------------
    # OCR 单张图片（Image 对象）
    # -----------------------
    @staticmethod
    def ocr_image_object(img: Image.Image) -> str:
        img_proc = OCRProcessor.preprocess_image(img)
        try:
            text = pytesseract.image_to_string(
                img_proc,
                lang=SystemConfig.OCR_LANG,
                timeout=SystemConfig.TESSERACT_TIMEOUT,
            )
        except Exception:
            text = ""

        text = text.strip()
        if text:
            return f"[IMAGE_OCR]\n{text}\n"
        return ""

    # -----------------------
    # OCR PNG/JPG 路径
    # -----------------------
    @staticmethod
    def ocr_image_file(path: str) -> str:
        try:
            img = Image.open(path)
        except Exception:
            return ""
        return OCRProcessor.ocr_image_object(img)

    # -----------------------
    # OCR PDF 中图片
    # -----------------------
    @staticmethod
    def extract_images_from_pdf(pdf_path: str) -> List[str]:
        """使用 pdf2image 将 PDF 转成图片，然后做 OCR"""
        texts = []
        try:
            pages = convert_from_path(pdf_path)
            for img in pages:
                txt = OCRProcessor.ocr_image_object(img)
                if txt.strip():
                    texts.append(txt)
        except Exception:
            pass
        return texts

    # -----------------------
    # OCR DOCX (word) 图片
    # -----------------------
    @staticmethod
    def extract_images_from_docx(docx_path: str) -> List[str]:
        texts = []
        try:
            d = docx.Document(docx_path)
            rels = d.part._rels
            for rel in rels:
                target = rels[rel]
                if "image" in target.target_ref:
                    img_bytes = target.target_part.blob
                    try:
                        img = Image.open(bytes(img_bytes))
                    except Exception:
                        continue

                    txt = OCRProcessor.ocr_image_object(img)
                    if txt.strip():
                        texts.append(txt)
        except Exception:
            pass
        return texts

    # -----------------------
    # OCR PPTX 图片
    # -----------------------
    @staticmethod
    def extract_images_from_pptx(pptx_path: str) -> List[str]:
        texts = []
        try:
            prs = Presentation(pptx_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # picture
                        img_blob = shape.image.blob
                        try:
                            img = Image.open(bytes(img_blob))
                        except Exception:
                            continue

                        txt = OCRProcessor.ocr_image_object(img)
                        if txt.strip():
                            texts.append(txt)
        except Exception:
            pass
        return texts


# ============================================================
# KnowledgeBaseManager：文档解析（文本 + 图片 OCR）
# ============================================================

class KnowledgeBaseManager:
    def __init__(self, vector_db: Chroma):
        self.vector_db = vector_db
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=SystemConfig.CHUNK_SIZE,
            chunk_overlap=SystemConfig.CHUNK_OVERLAP,
        )

    # -----------------------
    # 文本垃圾过滤
    # -----------------------
    @staticmethod
    def is_noise(text: str) -> bool:
        if not text or not text.strip():
            return True

        t = text.strip()

        # 1. 特殊符号占比 > 50%
        non_alpha_ratio = sum(
            1 for c in t if not (c.isalnum() or c in " \n，。,.!?？！")
        )
        if non_alpha_ratio / max(1, len(t)) > 0.5:
            return True

        # 2. 重复字符 (如 "aaaaaaa", "！！！！！")
        if re.search(r"(.)\1{7,}", t):
            return True

        # 3. 提示注入攻击
        injection_patterns = [
            "ignore previous instructions",
            "forget all previous instructions",
            "system prompt",
            "jailbreak",
            "you are chatgpt",
            "你现在的身份是",
            "你不是AI",
            "越狱",
            "system:",
        ]
        low = t.lower()
        if any(p in low for p in injection_patterns):
            return True

        # 4. 脏词过滤
        bad_words = {"fuck", "shit", "bitch", "操你", "傻逼", "sb", "妈的"}
        if any(b in low for b in bad_words):
            return True

        # 5. 纯 URL / 纯代码片段
        if len(t) < 20 and ("http://" in t or "https://" in t):
            return True

        return False

    # -----------------------
    # 读取 PDF 文本 + 图片 OCR
    # -----------------------
    def load_pdf(self, path: str) -> List[Document]:
        docs = []

        # 文本部分
        try:
            loader = PyPDFLoader(path)
            docs.extend(loader.load())
        except Exception as e:
            print("[RAG] PDF 文本解析失败:", e)

        # 图片 OCR 部分
        try:
            ocr_texts = OCRProcessor.extract_images_from_pdf(path)
            for t in ocr_texts:
                docs.append(Document(page_content=t, metadata={"ocr": True}))
        except Exception as e:
            print("[RAG] PDF OCR 失败:", e)

        return docs

    # -----------------------
    # 读取 DOCX 文本 + 图片 OCR
    # -----------------------
    def load_docx(self, path: str) -> List[Document]:
        docs = []

        # 文本
        try:
            d = docx.Document(path)
            text = "\n".join(
                p.text.strip() for p in d.paragraphs if p.text.strip()
            )
            if text.strip():
                docs.append(Document(page_content=text))
        except Exception as e:
            print("[RAG] DOCX 文本解析失败:", e)

        # 图片 OCR
        try:
            ocr = OCRProcessor.extract_images_from_docx(path)
            for t in ocr:
                docs.append(Document(page_content=t, metadata={"ocr": True}))
        except Exception as e:
            print("[RAG] DOCX OCR 失败:", e)

        return docs

    # -----------------------
    # 读取 PPTX 文本 + 图片 OCR
    # -----------------------
    def load_pptx(self, path: str) -> List[Document]:
        docs = []

        # 文本
        try:
            prs = Presentation(path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            if lines:
                docs.append(Document(page_content="\n".join(lines)))
        except Exception as e:
            print("[RAG] PPTX 文本解析失败:", e)

        # 图片 OCR
        try:
            ocr = OCRProcessor.extract_images_from_pptx(path)
            for t in ocr:
                docs.append(Document(page_content=t, metadata={"ocr": True}))
        except Exception as e:
            print("[RAG] PPTX OCR 失败:", e)

        return docs

    # -----------------------
    # 根据扩展名选择加载器
    # -----------------------
    def load_document(self, file_path: str) -> List[Document]:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self.load_pdf(file_path)
        elif ext == ".docx":
            return self.load_docx(file_path)
        elif ext == ".pptx":
            return self.load_pptx(file_path)
        else:
            raise ValueError(f"不支持的文件类型：{ext}")

    # -----------------------
    # 基础版上传（V1）—— V2 中会重写
    # -----------------------
    def upload_data(
        self,
        doc_id: str,
        file_path: str,
        original_name: Optional[str] = None,
        user_id: Optional[str] = None,
        subject: Optional[str] = None,
        unit: Optional[str] = None,
    ):
        print(f"[RAG] 开始解析文档：{file_path}")

        try:
            raw_docs = self.load_document(file_path)
        except Exception as e:
            print("[RAG] 文档加载失败:", e)
            return []

        # 删除旧向量
        try:
            self.vector_db.delete(where={"doc_id": doc_id})
        except Exception:
            pass

        # 分 chunk
        chunks = self.text_splitter.split_documents(raw_docs)

        valid_chunks = []
        subject_norm = normalize_subject_label(subject)

        for i, c in enumerate(chunks):
            if self.is_noise(c.page_content):
                continue

            c.metadata["doc_id"] = doc_id
            c.metadata["chunk_index"] = i
            if original_name:
                c.metadata["file_name"] = original_name
            if user_id:
                c.metadata["user_id"] = user_id
            if subject_norm:
                c.metadata["subject"] = subject_norm
            if unit:
                c.metadata["unit"] = unit

            valid_chunks.append(c)

        if not valid_chunks:
            print("[RAG] 所有 chunk 被过滤（可能是垃圾/注入）")
            return []

        self.vector_db.add_documents(valid_chunks)

        print(
            f"[RAG] 文档 {doc_id} 完成解析，共写入 {len(valid_chunks)} 个 chunks"
        )
        return [c.page_content for c in valid_chunks]


# ============================================================
# IntelligentAssistant：Router / Retrieval / Answer
# ============================================================

class IntelligentAssistant:
    def __init__(self, vector_db: Chroma, llm: ChatOpenAI):
        self.vector_db = vector_db
        self.llm = llm

        # 检索参数
        self.RECALL_K = 50            # 初次召回数量
        self.RERANK_TOP_K = 3         # 最终使用 top-k 文档
        # 为了兼顾中文 / 多语言检索，这里再放宽一些阈值
        self.MIN_OVERLAP_RATIO = 0.02 # 词面重叠最低比例（更宽松）
        self.STOPWORDS = set(GLOBAL_STOPWORDS)

    # ========================================================
    # 规则 Router：判断是否需要检索（非大模型部分）
    # ========================================================
    def _rule_router(self, query: str) -> Optional[bool]:
        """
        优先根据规则判定是否需要检索：
        - 若明显是闲聊 → False
        - 若明显问教材内容 / 项目要求 → True
        - 否则 → None（交给模型判定）
        """

        q = query.lower()

        # 英文 token（避免 "hi" 匹配到 "which" 的 bug）
        tokens = self._tokenize(query)

        # -------------------------
        # 1) 判断是否是闲聊（无需 RAG）
        #    - 英文：按 token 精确匹配
        #    - 中文：按短语包含匹配
        # -------------------------
        small_talk_words_en = ["hello", "hi"]
        small_talk_phrases_zh = [
            "哈哈",
            "你好",
            "早上好",
            "晚上好",
            "你是谁",
            "介绍一下",
            "心情",
            "天气",
        ]

        if any(w in tokens for w in small_talk_words_en):
            return False

        if any(w in q for w in small_talk_phrases_zh):
            return False

        # -------------------------
        # 2) 与教材 / 项目相关的关键字 → 强制使用 RAG
        # -------------------------
        rag_keywords_en = [
            "project a",
            "project requirement",
            "project requirements",
            "assignment",
            "deadline",
            "due date",
            "tech stack",
            "technology",
            "technologies",
            "course project",
            "ooad",
            "rag module",
        ]
        rag_keywords_zh = [
            "这页",
            "这一段",
            "这张图",
            "教材",
            "课本",
            "slides",
            "pdf",
            "请解释图",
            "第几页",
            "内容是什么",
            "项目要求",
            "课程项目",
            "作业要求",
            "技术栈",
            "需要哪些技术",
        ]

        if any(k in q for k in rag_keywords_en) or any(k in q for k in rag_keywords_zh):
            return True

        # -------------------------
        # 3) 注入攻击强制开启 RAG（交给教材兜底）
        # -------------------------
        jailbreak_patterns = [
            "ignore previous",
            "system prompt",
            "越狱",
            "jailbreak",
            "你现在的身份是",
            "你不是 ai",
        ]

        if any(w in q for w in jailbreak_patterns):
            return True

        return None  # 交给模型判定

    # ========================================================
    # 模型 Router：调用小模型判断是否需要检索
    # ========================================================
    def _model_router(self, query: str) -> bool:
        tpl = PromptTemplate.from_template(
            """You are a RAG router for a course assistant.
Your task: determine if answering this question requires the uploaded course materials (PDF, DOCX, PPTX, Markdown, or images).

If the user asks about project requirements, deadlines, grading rules, definitions, examples, theorem statements,
or anything likely in the course documents, answer YES.

If it is pure chit-chat or personal questions not related to the course,
answer NO.

You may receive questions in English or Chinese.
对于与课程教材、作业、Project A 相关的问题，请回答 YES。
对于纯闲聊、与课程无关的问题，请回答 NO。

Return ONLY: YES or NO.

Question: {q}
"""
        )
        try:
            res = (tpl | self.llm).invoke({"q": query})
            s = res.content.strip().upper()
            if "YES" in s and "NO" not in s:
                return True
            if "NO" in s and "YES" not in s:
                return False
        except Exception:
            # 模型失效 → 默认开启检索
            return True

        # 模糊情况 → 默认开启检索
        return True

    # ========================================================
    # 综合 Router（规则层 + 模型层）
    # ========================================================
    def _check_retrieval_necessity(self, query: str) -> bool:
        """规则判断在前，模型判断在后"""

        # 1) 先规则判断
        rule_result = self._rule_router(query)
        if rule_result is not None:
            print(f"[RAG] Router（规则层）= {rule_result}")
            return rule_result

        # 2) 再模型判断
        model_result = self._model_router(query)
        print(f"[RAG] Router（模型层）= {model_result}")
        return model_result

    # ========================================================
    # Tokenize（用于词面重叠和 Router 英文精确匹配）
    # 兼容：英文按单词、中文按单字 / CJK 字符
    # ========================================================
    def _tokenize(self, text: str) -> set:
        tokens = set()
        text = text.lower()

        # 1) 英文/数字按空格分词
        rough_tokens = re.split(r"\s+", text)
        for t in rough_tokens:
            t = t.strip('.,!?;:"\'()[]{}，。、？！')
            if not t:
                continue
            # 只要含有拉丁字母或数字，当作英文 token
            if re.search(r"[a-z0-9]", t):
                if t in self.STOPWORDS:
                    continue
                tokens.add(t)

        # 2) 中文 / 日文 / 其他 CJK 单字符 token
        for ch in text:
            # CJK Unified Ideographs 范围（覆盖中日文大部分汉字）
            if "\u4e00" <= ch <= "\u9fff":
                if ch in self.STOPWORDS:
                    continue
                tokens.add(ch)

        return tokens

    # ========================================================
    # 词面重叠比例计算：用于判断检索结果是否相关
    # ========================================================
    def _lexical_overlap(self, query: str, doc_text: str) -> float:
        q_tokens = self._tokenize(query)
        if not q_tokens:
            return 0.0
        d_tokens = self._tokenize(doc_text)
        if not d_tokens:
            return 0.0
        inter = q_tokens.intersection(d_tokens)
        return len(inter) / len(q_tokens)

    # ========================================================
    # Retrieval：从 Chroma 检索 + 词重叠过滤 + rerank
    # 注意：现在返回的是 Document 列表，方便后面拿到 metadata 做“来源段落”展示
    # 增强：加入 subject 过滤（user + subject 粒度的向量库视图）
    # ========================================================
    def _retrieve(
        self,
        query: str,
        user_id: Optional[str] = None,
        subject: Optional[str] = None,
    ) -> List[Document]:
        print(f"[RAG] 开始召回（k={self.RECALL_K}，user={user_id}, subject={subject})...")

        subject_norm = normalize_subject_label(subject)
        base_filter = {}
        if user_id:
            base_filter["user_id"] = user_id
        if subject_norm and subject_norm != "其他":
            base_filter["subject"] = subject_norm

        def _do_search(filter_dict):
            try:
                return self.vector_db.similarity_search_with_score(
                    query,
                    k=self.RECALL_K,
                    filter=filter_dict or None,
                )
            except Exception as e:
                print("[RAG] VectorStore 检索失败：", e)
                return []

        # 第一次：限定 user + subject
        docs_scores = _do_search(base_filter)

        # 若按学科检索为空，退回仅 user 过滤的全学科检索
        if not docs_scores and "subject" in base_filter:
            print("[RAG] 当前学科未召回到任何文档，退回到不带 subject 过滤的检索")
            base_filter.pop("subject", None)
            docs_scores = _do_search(base_filter)

        if not docs_scores:
            print("[RAG] 未召回到任何文档")
            return []

        scored = []  # (doc, distance, overlap)

        for doc, distance in docs_scores:
            overlap = self._lexical_overlap(query, doc.page_content)
            scored.append((doc, distance, overlap))

        # 展示前几条调试信息
        print("[RAG] 检索前 5 条：")
        for i, (doc, dis, ov) in enumerate(scored[:5]):
            print(f"  #{i + 1}: distance={dis:.4f}, overlap={ov:.4f}")

        # -------------------------
        # 过滤：保留 overlap >= MIN_OVERLAP_RATIO 的文档
        # -------------------------
        filtered = [
            (doc, dis, ov)
            for doc, dis, ov in scored
            if ov >= self.MIN_OVERLAP_RATIO
        ]

        if not filtered:
            print("[RAG] 词面重叠不足，认定为知识盲区")
            return []

        # -------------------------
        # reranking：按距离由小到大排序
        # -------------------------
        filtered.sort(key=lambda x: x[1])  # distance 越小越相关

        # 取 top-K
        top_docs_triplets = filtered[: self.RERANK_TOP_K]
        top_docs = [doc for doc, _, _ in top_docs_triplets]

        print(f"[RAG] 过滤后剩 {len(filtered)} 条，最终采用 {len(top_docs)} 条")

        return top_docs

    # ========================================================
    # GPT 回答模块（含 Timeout Fallback + 图片 OCR 支持）
    # ========================================================
    def _answer_with_gpt(
        self,
        query: str,
        chunks: List[str],
        image_chunks: Optional[List[str]] = None,
        subject: Optional[str] = None,
        unit: Optional[str] = None,
    ) -> str:
        """
        调用 GPT（通过代理）生成最终答案。

        需求：
        - 输出中必须包含：
          【问题分类】(subject / unit / has_support)
          【解题思路（Step-by-Step）】（无论有没有教材支撑都要有）
          【最终答案】
          【来源段落】（有支撑时直接贴 chunk 文本而不是序号）
        """

        from openai import OpenAI

        client = OpenAI(
            api_key=OPENAI_API_KEY,
            base_url=BASE_URL,
        )

        # 上下文状态
        image_chunks = image_chunks or []
        has_doc_chunks = bool(chunks)
        has_image_chunks = bool(image_chunks)

        subject_norm = normalize_subject_label(subject) or "其他"
        unit_norm = (unit or "通用").strip() or "通用"

        # “教材支撑”仅指来自知识库的文本 chunks
        has_support = has_doc_chunks
        has_support_str = "是" if has_support else "否"

        # -------------------------
        # 系统 prompt：约束输出结构 + 禁止幻觉
        # -------------------------
        system_prompt = f"""
你是一个严格受“课程教材知识库”约束的学习助手（RAG 模型）。

你会看到两类上下文：
1. 来自课程教材 / 项目的文本片段（Chunks）—— 这些属于“教材支撑”
2. 来自用户“本次提问”上传图片的 OCR 文本（ImageChunks）—— 属于临时上下文

此外，系统会提供当前问题所属的【学科】和【单元】信息：
- 学科(subject)：从 {SUBJECT_LABELS} 中的一个，例如 数 / 物 / 化 / 生 / 政 / 史 / 地 / 计算机 / 其他
- 单元(unit)：更细的章节或知识点名称，如“函数与导数”、“细胞结构”等

【关于 has_support 标记】：
- has_support = "是" 表示本次检索到至少一个教材 chunk（Chunks 非空）；
- has_support = "否" 表示本次没有检索到任何教材 chunk，此时：
  - 你仍然可以参考 ImageChunks（若存在），但它们不视为课程教材支撑；
  - 你必须在答案中明确说明“根据当前教材内容没有足够信息”。

【特别重要的规则】：
1. 你只能使用本次请求提供的 Chunks / ImageChunks 中的信息进行回答，
   **绝不能**使用任何外部世界知识或常识来补全结论。
2. 无论 has_support 是“是”还是“否”，你都必须给出【解题思路（Step-by-Step）】。
3. 如果 has_support = "是"，在【来源段落】中要**直接粘贴你真正用到的 chunk 文本**，
   可以适当精简，但必须包含关键句子，而不是只写“Chunk 1/2/3”等编号。
4. 如果 has_support = "否"，在【最终答案】中必须明确指出：
   “根据当前教材内容，没有足够信息得出确切结论。”，
   在【来源段落】中写：
   “无（当前知识库中未检索到相关教材内容）”。
5. 如果只有 ImageChunks（没有教材 chunks），你可以在【解题思路】中说明
   你是根据图片 OCR 文本做出的判断，同时仍需说明这不属于正式教材支撑。

【输出格式要求】—— 必须严格按照下面模板输出，不要增加或删除一级标题：

【问题分类】
- 学科(subject)：...
- 单元(unit)：...
- 是否有教材支撑(has_support)：是/否

【解题思路（Step-by-Step）】
1. ...
2. ...
（必要时可以 3., 4. 等）

【最终答案】
...

【来源段落】
- 若 has_support = "是"：逐条粘贴你用到的教材内容，例如：
  - Chunk 1：<对应 chunk 的关键句子或全文>
  - Chunk 2：<...>
- 若 has_support = "否"：写“无（当前知识库中未检索到相关教材内容）”
"""

        # -------------------------
        # 构造用户 prompt（包含分类 + has_support 标记 + 上下文）
        # -------------------------
        user_prompt = ""
        user_prompt += f"【问题学科分类（供你参考）】subject={subject_norm}, unit={unit_norm}\n"
        user_prompt += f"【是否有教材支撑】has_support={has_support_str}\n\n"
        user_prompt += f"用户问题：{query}\n\n"

        if has_doc_chunks:
            user_prompt += "以下是从课程教材 / 项目知识库中检索到的相关内容片段（Chunks）：\n\n"
            for i, c in enumerate(chunks):
                user_prompt += f"[Chunk {i + 1}]\n{c}\n\n"

        if has_image_chunks:
            user_prompt += "以下是来自用户本次上传图片的 OCR 文本（ImageChunks）：\n\n"
            for i, c in enumerate(image_chunks):
                user_prompt += f"[ImageChunk {i + 1}]\n{c}\n\n"

        # -------------------------
        # 调用 GPT（使用新版 openai SDK 接口）
        # -------------------------
        try:
            start = time.time()

            response = client.chat.completions.create(
                model=RAG_ANSWER_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
            )

            # 若 GPT 生成过慢（自定义判断）
            if time.time() - start > SystemConfig.RAG_TIMEOUT:
                raise TimeoutError("GPT timeout exceeded")

            return response.choices[0].message.content.strip()

        except Exception as e:
            print("[RAG] GPT 调用失败：", e)

            # Timeout / 其它异常 fallback：保守回答
            summary_parts = []
            for c in chunks[:2]:
                summary_parts.append(f"- [教材片段] {c[:150]}...")
            for c in image_chunks[:2]:
                summary_parts.append(f"- [图片片段] {c[:150]}...")

            if not summary_parts:
                return "由于系统繁忙或网络错误，目前无法生成回答。"

            return (
                "由于系统繁忙或网络错误，目前无法生成完整回答。\n"
                "但根据已检索到的部分内容，你可以参考以下信息：\n\n"
                + "\n".join(summary_parts)
            )

    # ========================================================
    # 主入口：多图 OCR → Router → Retrieval → Answer（基础版）
    # ========================================================
    def handle_user_query(
        self,
        user_id: str,
        query: str,
        image_paths: Optional[List[str]] = None,
    ):
        print(f"[RAG] 用户 {user_id} 提问：{query}")

        # Step 1: OCR 多张图片
        ocr_text = ""
        image_chunks: List[str] = []

        if image_paths:
            print(f"[RAG] 共接收 {len(image_paths)} 张图片，开始 OCR ...")
            for p in image_paths:
                try:
                    t = OCRProcessor.ocr_image_file(p)
                    if t.strip():
                        ocr_text += t + "\n"
                        image_chunks.append(t.strip())
                except Exception as e:
                    print(f"[RAG] OCR 失败 ({p})：", e)

        full_query = (query + "\n" + ocr_text).strip()
        print("[RAG] 完整 Query：", full_query)

        # Step 2: Router
        need_retrieval = self._check_retrieval_necessity(full_query)
        print(f"[RAG] 是否需要检索：{need_retrieval}")

        chunks: List[str] = []
        retrieved_docs: List[Document] = []
        subject = None
        unit = None

        if need_retrieval:
            retrieved_docs = self._retrieve(full_query, user_id=user_id, subject=None)
            chunks = [d.page_content for d in retrieved_docs]

        # Step 4: Answer
        final_answer = self._answer_with_gpt(
            full_query,
            chunks,
            image_chunks=image_chunks,
            subject=subject,
            unit=unit,
        )

        # Step 5: 组装来源信息
        source_chunks = []
        for d in retrieved_docs:
            source_chunks.append(
                {
                    "text": d.page_content,
                    "metadata": getattr(d, "metadata", {}),
                }
            )

        return {
            "query": full_query,
            "retrieval_performed": need_retrieval,
            "matched_chunks": chunks,
            "source_chunks": source_chunks,
            "subject": subject,
            "unit": unit,
            "final_answer": final_answer,
        }


# ============================================================
# 安全工具 / OCR v2 / NoiseFilter v2 / Logger
# ============================================================

class SecurityUtils:

    JAILBREAK_PATTERNS = [
        r"ignore (all )?previous instructions",
        r"forget (all )?previous instructions",
        r"jailbreak",
        r"越狱",
        r"system prompt",
        r"pretend to be",
        r"do anything now",
        r"伪装成",
        r"你不再是",
        r"你现在的身份是",
        r"你是一个人类",
        r"你不是一个 AI",
        r"绕过规则",
        r"请忽略所有内容",
        r"重置你的身份为",
        r"你在扮演",
    ]

    @staticmethod
    def is_prompt_injection(text: str) -> bool:
        low = text.lower()
        for pattern in SecurityUtils.JAILBREAK_PATTERNS:
            if re.search(pattern, low):
                return True
        return False

    @staticmethod
    def clean_text(text: str) -> str:
        if not text:
            return ""

        # 零宽字符
        text = re.sub(r"[\u200b\u200c\u200d\ufeff]", "", text)

        # 隐藏控制字符
        text = "".join(
            c for c in text if c.isprintable() or c in " \n\t，。！？"
        )

        # Emoji（替换为空）
        emoji_pattern = re.compile(
            "["
            "\U0001F600-\U0001F64F"
            "\U0001F300-\U0001F5FF"
            "\U0001F680-\U0001F6FF"
            "\U0001F1E0-\U0001F1FF"
            "]+",
            flags=re.UNICODE,
        )
        text = emoji_pattern.sub("", text)

        return text.strip()


class OCRProcessorV2(OCRProcessor):

    MAX_OCR_PIXELS = 4_000_000  # 超过此像素自动缩放（2000x2000）

    @staticmethod
    def safe_resize(img: Image.Image) -> Image.Image:
        w, h = img.size
        if w * h > OCRProcessorV2.MAX_OCR_PIXELS:
            factor = (OCRProcessorV2.MAX_OCR_PIXELS / (w * h)) ** 0.5
            new_w = int(w * factor)
            new_h = int(h * factor)
            img = img.resize((new_w, new_h), Image.LANCZOS)
        return img

    @staticmethod
    def ocr_image_object(img: Image.Image) -> str:
        try:
            img = OCRProcessorV2.safe_resize(img)
        except Exception:
            pass
        return super(OCRProcessorV2, OCRProcessorV2).ocr_image_object(img)


class NoiseFilter:

    @staticmethod
    def is_noise(text: str) -> bool:
        if not text or not text.strip():
            return True

        t = text.strip()

        # 1. 特殊字符占比 > 60%
        non_alpha_ratio = sum(
            1 for c in t if not (c.isalnum() or c in " \n\t，。,.!?？！")
        )
        if non_alpha_ratio / max(1, len(t)) > 0.6:
            return True

        # 2. 大量重复
        if re.search(r"(.)\1{6,}", t):
            return True

        # 3. 乱码（无正常单词/中文）
        if not re.search(r"[A-Za-z0-9\u4e00-\u9fa5]", t):
            return True

        # 4. 只是 PPT 页码/标题（内容太短）
        short_patterns = [
            r"^第.{0,3}页$",
            r"^谢谢$",
            r"^目录$",
            r"^封面$",
        ]
        for sp in short_patterns:
            if re.match(sp, t):
                return True

        # 5. 注入直接过滤
        if SecurityUtils.is_prompt_injection(t):
            return True

        return False


class RAGLogger:
    @staticmethod
    def log(*args):
        print("[RAG LOG]", *args)

    @staticmethod
    def warn(*args):
        print("[RAG WARN]", *args)

    @staticmethod
    def error(*args):
        print("[RAG ERROR]", *args)


# ============================================================
# KnowledgeBaseManagerV2：集成 Parser / OCR / 安全过滤
# ============================================================

class KnowledgeBaseManagerV2(KnowledgeBaseManager):

    MAX_PDF_PAGES = 250
    MAX_DOC_SIZE = 30 * 1024 * 1024  # 最大 30MB

    def __init__(self, vector_db: Chroma):
        super().__init__(vector_db)
        try:
            self.parser = Parser()
        except Exception:
            self.parser = None
            RAGLogger.warn("Parser 初始化失败，将仅使用旧版解析逻辑")

    def _parse_with_parser(self, file_path: str, ext: str) -> List[Document]:
        if self.parser is None:
            return []

        try:
            res = self.parser.parse_file(
                file_path,
                ocr=True,
                ocr_lang=SystemConfig.OCR_LANG,
            )
        except Exception as e:
            RAGLogger.warn("Parser 解析失败：", e)
            return []

        docs: List[Document] = []

        pages = res.get("pages") if isinstance(res, dict) else None
        if pages and isinstance(pages, list):
            for idx, page_text in enumerate(pages):
                if not page_text:
                    continue
                text = str(page_text).strip()
                if not text:
                    continue
                docs.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": str(file_path),
                            "page": idx + 1,
                            "parser": True,
                        },
                    )
                )
        else:
            text = ""
            if isinstance(res, dict):
                text = str(res.get("text", "")).strip()
            else:
                text = str(res).strip()
            if text:
                docs.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": str(file_path),
                            "parser": True,
                        },
                    )
                )

        return docs

    @staticmethod
    def _check_file_size(path: str):
        size = os.path.getsize(path)
        if size > KnowledgeBaseManagerV2.MAX_DOC_SIZE:
            raise ValueError(
                f"文件过大（>{KnowledgeBaseManagerV2.MAX_DOC_SIZE} bytes），拒绝解析"
            )

    def load_pdf(self, path: str) -> List[Document]:
        self._check_file_size(path)

        docs = []

        try:
            loader = PyPDFLoader(path)
            raw_docs = loader.load()
            if len(raw_docs) > KnowledgeBaseManagerV2.MAX_PDF_PAGES:
                raw_docs = raw_docs[: KnowledgeBaseManagerV2.MAX_PDF_PAGES]
            docs.extend(raw_docs)
        except Exception as e:
            RAGLogger.warn("PDF 文本解析失败：", e)

        try:
            ocr_texts = OCRProcessorV2.extract_images_from_pdf(path)
            for t in ocr_texts:
                docs.append(
                    Document(
                        page_content=t,
                        metadata={"ocr": True, "source": "pdf_image"},
                    )
                )
        except Exception as e:
            RAGLogger.warn("PDF OCR 失败：", e)

        return docs

    def load_docx(self, path: str) -> List[Document]:
        self._check_file_size(path)
        docs = []

        try:
            d = docx.Document(path)
            text = "\n".join(
                p.text.strip() for p in d.paragraphs if p.text.strip()
            )
            if text.strip():
                docs.append(
                    Document(
                        page_content=text, metadata={"source": "docx_text"}
                    )
                )
        except Exception as e:
            RAGLogger.warn("DOCX 文本解析失败：", e)

        try:
            ocr_texts = OCRProcessorV2.extract_images_from_docx(path)
            for t in ocr_texts:
                docs.append(
                    Document(
                        page_content=t,
                        metadata={"ocr": True, "source": "docx_image"},
                    )
                )
        except Exception as e:
            RAGLogger.warn("DOCX OCR 失败：", e)

        return docs

    def load_pptx(self, path: str) -> List[Document]:
        self._check_file_size(path)
        docs = []

        try:
            prs = Presentation(path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        cleaned = SecurityUtils.clean_text(
                            shape.text.strip()
                        )
                        if cleaned:
                            lines.append(cleaned)
            if lines:
                docs.append(
                    Document(
                        page_content="\n".join(lines),
                        metadata={"source": "pptx_text"},
                    )
                )
        except Exception as e:
            RAGLogger.warn("PPTX 文本解析失败：", e)

        try:
            ocr = OCRProcessorV2.extract_images_from_pptx(path)
            for t in ocr:
                docs.append(
                    Document(
                        page_content=t,
                        metadata={"ocr": True, "source": "pptx_image"},
                    )
                )
        except Exception as e:
            RAGLogger.warn("PPTX OCR 失败：", e)

        return docs

    def upload_data(
        self,
        doc_id: str,
        file_path: str,
        original_name: Optional[str] = None,
        user_id: Optional[str] = None,
        subject: Optional[str] = None,
        unit: Optional[str] = None,
    ):
        RAGLogger.log("开始解析文档：", file_path)

        ext = os.path.splitext(file_path)[1].lower()

        if ext in (".pdf", ".docx", ".pptx"):
            self._check_file_size(file_path)

        raw_docs: List[Document] = []

        if hasattr(self, "parser") and self.parser is not None:
            raw_docs = self._parse_with_parser(file_path, ext)

        if not raw_docs:
            if ext == ".pdf":
                raw_docs = self.load_pdf(file_path)
            elif ext == ".docx":
                raw_docs = self.load_docx(file_path)
            elif ext == ".pptx":
                raw_docs = self.load_pptx(file_path)
            elif ext in (".md", ".markdown"):
                try:
                    with open(
                        file_path,
                        "r",
                        encoding="utf-8",
                        errors="ignore",
                    ) as f:
                        text = f.read()
                    if text.strip():
                        raw_docs = [
                            Document(
                                page_content=text,
                                metadata={"source": "markdown_text"},
                            )
                        ]
                except Exception as e:
                    RAGLogger.warn("Markdown 文本解析失败：", e)
            elif ext in (
                ".png",
                ".jpg",
                ".jpeg",
                ".bmp",
                ".gif",
                ".webp",
                ".tif",
                ".tiff",
            ):
                try:
                    img_text = OCRProcessorV2.ocr_image_file(file_path)
                    if img_text.strip():
                        raw_docs = [
                            Document(
                                page_content=img_text,
                                metadata={"ocr": True, "source": "image"},
                            )
                        ]
                except Exception as e:
                    RAGLogger.warn("图片 OCR 失败：", e)
            else:
                raise ValueError(f"不支持的文件类型：{ext}")

        try:
            self.vector_db.delete(where={"doc_id": doc_id})
        except Exception:
            pass

        chunks = self.text_splitter.split_documents(raw_docs)

        valid = []
        subject_norm = normalize_subject_label(subject)
        unit_norm = (unit or "").strip()

        for i, c in enumerate(chunks):
            text = SecurityUtils.clean_text(c.page_content)

            if NoiseFilter.is_noise(text):
                continue

            if SecurityUtils.is_prompt_injection(text):
                continue

            c.page_content = text
            c.metadata["doc_id"] = doc_id
            c.metadata["file_id"] = doc_id
            c.metadata["chunk_index"] = i
            if original_name:
                c.metadata["file_name"] = original_name
            if user_id:
                c.metadata["user_id"] = user_id
            if subject_norm:
                c.metadata["subject"] = subject_norm
            if unit_norm:
                c.metadata["unit"] = unit_norm

            valid.append(c)

        if not valid:
            RAGLogger.warn(
                "所有 chunk 已被判定为垃圾或注入，未写入数据库"
            )
            return []

        self.vector_db.add_documents(valid)
        RAGLogger.log(f"文档 {doc_id} 已成功写入 {len(valid)} 个 chunks")

        return [c.page_content for c in valid]


# ============================================================
# IntelligentAssistantV2：增加学科/单元路由
# ============================================================

class IntelligentAssistantV2(IntelligentAssistant):

    def _clean_query(self, q: str) -> str:
        q = SecurityUtils.clean_text(q)
        return q

    def _rule_subject(self, text: str) -> Optional[str]:
        if not text:
            return None
        lower = text.lower()
        scores = {label: 0 for label in SUBJECT_LABELS}

        for label, kws in SUBJECT_KEYWORDS.items():
            for kw in kws:
                if not kw:
                    continue
                if kw.lower() in lower:
                    scores[label] += 1

        best_label = max(scores, key=scores.get)
        if scores[best_label] == 0:
            return None

        if best_label == "其他" and scores[best_label] < 2:
            return None

        return best_label

    def _llm_subject_and_unit(self, query: str) -> Tuple[str, str]:
        tpl = PromptTemplate.from_template(
            """
你是一个“学科与单元分类器”，需要根据学生的问题或指令，判断它属于哪一个学科，并给出一个更细的“单元/章节”名称。

可选学科标签（必须严格从中选择一个返回到 JSON 的 subject 字段）：
- 数：数学相关（代数、几何、函数、微积分、概率统计等）
- 物：物理相关（力学、电学、光学、热学、近代物理等）
- 化：化学相关（无机/有机化学、化学方程式、酸碱、氧化还原等）
- 生：生物相关（细胞、遗传、生态、进化、生理等）
- 政：政治/思政/政治经济学相关（国家、公民、法律、市场经济、宏观调控等）
- 史：历史相关（中国史、世界史、古代史、近现代史等）
- 地：地理相关（自然地理、人文地理、区域地理、经济地理等）
- 计算机：计算机与信息技术（编程、算法、数据结构、操作系统、网络、数据库等）
- 其他：不属于以上任何一类，或者跨学科综合、泛化问题

请用 JSON 格式输出，字段为：
- subject: 上述学科标签之一（数/物/化/生/政/史/地/计算机/其他）
- unit: 一个简短的中文名称，描述该问题更细的单元/章节（如 “函数与导数”、“细胞结构与功能”）。
        若无法判断，请返回 "通用"。

注意：
- 不要输出多余的文字，不要加解释。
- 如确实不确定学科，请将 subject 设为 "其他"。

示例输出：
{{"subject": "数", "unit": "函数与导数"}}

待分类的问题：
{q}
"""
        )
        subject = "其他"
        unit = "通用"
        try:
            res = (tpl | self.llm).invoke({"q": query})
            raw = res.content.strip()

            raw = re.sub(r"^```json", "", raw, flags=re.IGNORECASE).strip()
            raw = re.sub(r"^```", "", raw).strip()
            raw = re.sub(r"```$", "", raw).strip()

            data = json.loads(raw)
            s = data.get("subject", "其他")
            u = data.get("unit", "通用")
            subject = normalize_subject_label(s) or "其他"
            unit = (u or "通用").strip() or "通用"
        except Exception as e:
            RAGLogger.warn("LLM 学科/单元路由失败，使用默认值：", e)
        return subject, unit

    def _classify_subject_and_unit(self, query: str) -> Tuple[str, str]:
        rule_subject = self._rule_subject(query)
        if rule_subject and rule_subject in SUBJECT_LABELS and rule_subject != "其他":
            return rule_subject, "通用"

        subject, unit = self._llm_subject_and_unit(query)
        subject = normalize_subject_label(subject) or "其他"
        unit = (unit or "通用").strip() or "通用"
        return subject, unit

    def handle_user_query(
        self,
        user_id: str,
        query: str,
        image_paths: Optional[List[str]] = None,
    ):
        RAGLogger.log(f"收到用户 {user_id} 查询：{query}")

        query = self._clean_query(query)

        ocr_text = ""
        image_chunks: List[str] = []

        if image_paths:
            RAGLogger.log(f"共 {len(image_paths)} 张图片，开始 OCR ...")
            for p in image_paths:
                try:
                    t = OCRProcessorV2.ocr_image_file(p)
                    if t.strip():
                        ocr_text += t + "\n"
                        image_chunks.append(t.strip())
                except Exception as e:
                    RAGLogger.error("OCR 错误：", e)

        full_query = (query + "\n" + ocr_text).strip()

        if SecurityUtils.is_prompt_injection(full_query):
            RAGLogger.warn("检测到提示注入攻击，已过滤危险内容")
            full_query = "用户试图攻击模型，请仅根据教材内容回答。"

        RAGLogger.log("最终 Query：", full_query)

        subject, unit = self._classify_subject_and_unit(full_query)
        RAGLogger.log(f"问题学科路由：subject={subject}, unit={unit}")

        need_rag = self._check_retrieval_necessity(full_query)
        RAGLogger.log("Router 结果：需要检索 =", need_rag)

        chunks: List[str] = []
        retrieved_docs: List[Document] = []
        if need_rag:
            retrieved_docs = self._retrieve(full_query, user_id=user_id, subject=subject)
            chunks = [d.page_content for d in retrieved_docs]

        # 是否有教材支撑（仅指知识库 chunks）
        has_support = bool(chunks)
        # 是否有图片 OCR 上下文
        has_image_context = bool(image_chunks)

        final_answer = self._answer_with_gpt(
            full_query,
            chunks,
            image_chunks=image_chunks,
            subject=subject,
            unit=unit,
        )

        source_chunks = []
        for d in retrieved_docs:
            source_chunks.append(
                {
                    "text": d.page_content,
                    "metadata": getattr(d, "metadata", {}),
                }
            )

        return {
            "query": full_query,
            "retrieval_performed": need_rag,
            "has_support": has_support,
            "has_image_context": has_image_context,
            "matched_chunks": chunks,
            "source_chunks": source_chunks,
            "subject": subject,
            "unit": unit,
            "final_answer": final_answer,
        }


# ============================================================
# 单例对象：V2 版本替换旧版本
# ============================================================

try:
    os.makedirs(SystemConfig.PERSIST_DIRECTORY, exist_ok=True)
except Exception as e:
    RAGLogger.error("初始化向量库目录失败：", e)

kb_manager = KnowledgeBaseManagerV2(vectorstore)
assistant = IntelligentAssistantV2(vectorstore, llm_router)


def _test_ocr_debug(image_path: str):
    try:
        txt = OCRProcessorV2.ocr_image_file(image_path)
        print("------ OCR DEBUG RESULT ------")
        print(txt)
        print("------ END ------")
    except Exception as e:
        print("[OCR DEBUG ERROR]", e)


def get_vectorstore_stats():
    try:
        stats = vectorstore._collection.count()
        print(f"[RAG] VectorStore 当前存储 chunks 数：{stats}")
        return stats
    except Exception as e:
        print("[RAG] 获取向量库统计失败：", e)
        return 0


RAGLogger.log(
    "RAG 服务初始化完成：OCR + 文本解析 + Parser 图文解析 + 多语言向量模型 + 学科路由 + 强检索 + 防注入 + 超时保护 + 图片直连大模型 已全部启用。"
)
